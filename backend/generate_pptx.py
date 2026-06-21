"""
Minimal Universität-Passau-style thesis outline deck (white background).

• White backgrounds; orange accent rule; grey Passau typography
• Raster wordmark (orange bar + Arial) — official logo hot-link returns HTTP 403
• Slide numbers bottom right:  n / N
• Every bullet ≤ 6 words (asserted at runtime)

Run:  python generate_pptx.py
Out:  EU_AI_Governance_Project_Overview.pptx
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PASSAU_ORANGE = RGBColor(0xFD, 0x89, 0x00)
PASSAU_GREY = RGBColor(0x4B, 0x4B, 0x4B)
PASSAU_LINE = RGBColor(0xC9, 0xC8, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

ASSETS_DIR = Path(__file__).resolve().parent / "assets"
WORDMARK = ASSETS_DIR / "uni_passau_wordmark.png"
OUTFILE = Path(__file__).resolve().parent / "EU_AI_Governance_Project_Overview.pptx"


def _word_count(s: str) -> int:
    cleaned = s.replace("·", " ").replace(";", " ").replace(":", " ").replace(",", " ")
    return len([w for w in cleaned.split() if w])


def _assert_bullets(items: list[str], max_w: int = 6) -> None:
    for t in items:
        n = _word_count(t)
        if n > max_w:
            raise ValueError(f"Bullet > {max_w} words ({n}): {t!r}")


def _ensure_wordmark_png() -> Path:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    if WORDMARK.exists():
        return WORDMARK
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        raise RuntimeError("pip install Pillow") from e

    wpx, hpx = 520, 72
    img = Image.new("RGBA", (wpx, hpx), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 12, 10, hpx - 12], fill=(253, 137, 0, 255))

    font = None
    for fp in (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
        "C:\\Windows\\Fonts\\arial.ttf",
    ):
        if os.path.isfile(fp):
            font = ImageFont.truetype(fp, 26)
            break
    if font is None:
        font = ImageFont.load_default()
    draw.text((22, 22), "Universität Passau", fill=(75, 75, 75, 255), font=font)
    img.save(WORDMARK, "PNG")
    return WORDMARK


def _bg_white(slide) -> None:
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = WHITE


def _rule_y(slide, y) -> None:
    sh = slide.shapes.add_shape(1, Inches(0), y, W, Inches(0.02))
    sh.fill.solid()
    sh.fill.fore_color.rgb = PASSAU_LINE
    sh.line.fill.background()


def _logo(slide, path: Path) -> None:
    slide.shapes.add_picture(str(path), W - Inches(2.35), Inches(0.14), width=Inches(2.15))


def _num(slide, i: int, n: int) -> None:
    tb = slide.shapes.add_textbox(W - Inches(1.05), H - Inches(0.48), Inches(0.9), Inches(0.4))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{i} / {n}"
    r.font.size = Pt(11)
    r.font.color.rgb = PASSAU_GREY


def _title(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.72), Inches(9.5), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = PASSAU_GREY
    ln = slide.shapes.add_shape(1, Inches(0.55), Inches(1.28), Inches(5.8), Inches(0.04))
    ln.fill.solid()
    ln.fill.fore_color.rgb = PASSAU_ORANGE
    ln.line.fill.background()


def _bullets(slide, items: list[str], l, t, w, h, sz=15) -> None:
    _assert_bullets(items, 6)
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for txt in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"•  {txt}"
        r.font.size = Pt(sz)
        r.font.color.rgb = PASSAU_GREY


def _toc(slide, l, t, w, h) -> None:
    rows: list[tuple[str, list[str]]] = [
        ("1.  Motivation & Problem Statement", ["EU AI Act as regulatory driver"]),
        ("2.  EU AI Act Key Concepts", ["Risk tiers, obligations, prohibited practices"]),
        ("3.  Project Goal & Vision", ["Finished product; live demo walkthrough"]),
        ("4.  System Architecture", ["Five-node pipeline; data flow"]),
        ("5.  Technical Deep-Dive", ["LlamaIndex · AutoGen · Pydantic AI"]),
        ("6.  Technical Constraints & Mitigations", ["Loops, Python 3.9, cost guards"]),
        ("7.  Conclusion & Future Work", ["Lessons learned; production roadmap"]),
    ]
    flat = [a for a, _ in rows] + [b for _, bb in rows for b in bb]
    _assert_bullets(flat, 6)

    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for head, subs in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = f"●  {head}"
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = PASSAU_GREY
        for sub in subs:
            p = tf.add_paragraph()
            p.level = 1
            p.space_after = Pt(3)
            r2 = p.add_run()
            r2.text = f"◦  {sub}"
            r2.font.size = Pt(14)
            r2.font.color.rgb = PASSAU_GREY


def _icons_vertical(slide, syms: list[str], cx, y0) -> None:
    y = y0
    for s in syms:
        tb = slide.shapes.add_textbox(cx - Inches(0.35), y, Inches(0.7), Inches(0.52))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = s
        r.font.size = Pt(26)
        y += Inches(0.58)


def build_deck() -> Path:
    logo = _ensure_wordmark_png()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slides_list: list = []

    def blank():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _bg_white(s)
        _logo(s, logo)
        _rule_y(s, Inches(1.18))
        slides_list.append(s)
        return s

    # ---- 1 title only ----
    s1 = blank()
    t1 = s1.shapes.add_textbox(Inches(0.85), Inches(2.55), Inches(11.6), Inches(0.95))
    p = t1.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Enterprise AI Governance Workflow Engine"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = PASSAU_GREY
    t2 = s1.shapes.add_textbox(Inches(0.85), Inches(4.0), Inches(11.6), Inches(0.55))
    p2 = t2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Mustafa Khalid"
    r2.font.size = Pt(22)
    r2.font.color.rgb = PASSAU_GREY

    # ---- 2 outline ----
    s2 = blank()
    _title(s2, "Presentation Outline")
    _toc(s2, Inches(0.55), Inches(1.48), Inches(8.0), Inches(5.45))
    _icons_vertical(s2, ["📑", "⚖️", "🎯", "🔧", "🧠", "⚠️", "🔭"], W - Inches(1.95), Inches(1.52))

    # ---- 3 motivation ----
    s3 = blank()
    _title(s3, "1.  Motivation & Problem Statement")
    _bullets(
        s3,
        [
            "AI Act drives strong compliance pressure",
            "Large fines threaten unacceptable risk systems",
            "Teams still lack early automated assessments",
            "This prototype closes that governance gap",
        ],
        Inches(0.55),
        Inches(1.45),
        Inches(7.85),
        Inches(5.2),
    )
    _icons_vertical(s3, ["📜", "⚡", "🔍", "✅"], W - Inches(1.95), Inches(1.55))

    # ---- 4 EU concepts ----
    s4 = blank()
    _title(s4, "2.  EU AI Act Key Concepts")
    _bullets(
        s4,
        [
            "Three risk tiers define legal obligations",
            "Certain AI practices stay fully prohibited",
            "High-risk systems demand robust safeguards",
            "Low tiers still owe transparency duties",
        ],
        Inches(0.55),
        Inches(1.45),
        Inches(7.85),
        Inches(5.2),
    )
    _icons_vertical(s4, ["🚫", "⚠️", "✅"], W - Inches(1.95), Inches(2.05))

    # ---- 5 goal / product / implementation ----
    s5 = blank()
    _title(s5, "3.  Project Goal & Vision")
    _bullets(
        s5,
        [
            "Automates EU AI Act compliance assessments",
            "Delivers validated JSON audit from descriptions",
            "LangGraph orchestrates every pipeline software step",
            "LlamaIndex AutoGen Pydantic AI Streamlit stack",
        ],
        Inches(0.55),
        Inches(1.45),
        Inches(7.85),
        Inches(5.2),
    )
    _icons_vertical(s5, ["🎯", "📋", "🧩", "💻"], W - Inches(1.95), Inches(1.5))

    # ---- 6 architecture ----
    s6 = blank()
    _title(s6, "4.  System Architecture")
    _bullets(
        s6,
        [
            "Ingest plain user system specification text",
            "Retrieve relevant law via LlamaIndex vectors",
            "Two AutoGen agents debate residual risks",
            "Schema step uses Pydantic AI validation",
            "Retry debate if validation remains failing",
            "Iteration cap prevents infinite LangGraph loops",
        ],
        Inches(0.55),
        Inches(1.45),
        Inches(7.85),
        Inches(5.2),
    )
    _icons_vertical(s6, ["➊", "➋", "➌", "➍", "↩️", "🛑"], W - Inches(1.95), Inches(1.35))

    # ---- 7 technical ----
    s7 = blank()
    _title(s7, "5.  Technical Deep-Dive")
    _bullets(
        s7,
        [
            "LlamaIndex embeds one local markdown corpus",
            "AutoGen runs architect versus compliance officer",
            "Pydantic AI enforces strict audit schema",
            "Streamlit hosts the seminar demonstration interface",
        ],
        Inches(0.55),
        Inches(1.45),
        Inches(7.85),
        Inches(5.2),
    )
    _icons_vertical(s7, ["📚", "🤝", "🔒", "🖥️"], W - Inches(1.95), Inches(1.55))

    # ---- 8 constraints ----
    s8 = blank()
    _title(s8, "6.  Technical Constraints & Mitigations")
    _bullets(
        s8,
        [
            "Pinned LlamaIndex supports Python three nine",
            "openai package pinned strictly below v2",
            "LLM usage implies ongoing API spend",
            "UserProxy stays outside debate agent pool",
        ],
        Inches(0.55),
        Inches(1.45),
        Inches(7.85),
        Inches(5.2),
    )
    _icons_vertical(s8, ["🐍", "📦", "💸", "🤝"], W - Inches(1.95), Inches(1.65))

    # ---- 9 conclusion ----
    s9 = blank()
    _title(s9, "7.  Conclusion & Future Work")
    _bullets(
        s9,
        [
            "Prototype proves full governance workflow path",
            "Adopt official EUR-Lex legal corpus next",
            "Add database logs and authentication next",
        ],
        Inches(0.55),
        Inches(1.45),
        Inches(7.85),
        Inches(5.2),
    )
    _icons_vertical(s9, ["✔️", "📚", "🗄️"], W - Inches(1.95), Inches(2.0))

    # ---- 10 thank you ----
    s10 = blank()
    c1 = s10.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(0.9))
    p = c1.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank you"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = PASSAU_GREY
    c2 = s10.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.7))
    p2 = c2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Any questions?"
    r2.font.size = Pt(26)
    r2.font.color.rgb = PASSAU_GREY

    total = len(slides_list)
    for idx, sl in enumerate(slides_list, start=1):
        _num(sl, idx, total)

    prs.save(str(OUTFILE))
    return OUTFILE


if __name__ == "__main__":
    p = build_deck()
    print(p)
